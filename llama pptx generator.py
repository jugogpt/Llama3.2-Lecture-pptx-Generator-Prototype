import requests
import json
import os
import ast
import subprocess
import numpy as np
import re
from pptx import Presentation
from pptx.util import Inches, Pt

# === LLaMA API Configuration ===
url = 'http://localhost:11434/api/generate'  # URL endpoint for local LLaMA 3.2 model via Ollama

# Prepare the prompt and configuration data for generating the PowerPoint blueprint
data = {
    'model': 'llama3.2',
    'parameter': "temperature 0.1",  # Low temperature for deterministic output
    'prompt': 'Here is your transcript to reference: {Transcript.txt}, here here is your powerpoint reference: {Template.pptx}' r"""Given this lecture transcript here: https://ocw.mit.edu/courses/7-012-introduction-to-biology-fall-2004/d79f25e7725465922f61117286158f6f_CovlKXmuWo.pdf, create a 20 slide powerpoint presentation, giving me the information for each of the slides in this format : give me the bullets for each slide in an array of strings. Have each array of strings be in a array of arrays, and past that in one line in python form and that is all you output. Here is the transcript and template, 

    create an array of arrays in this structure:

    1. the outer array is an array where each element is an array that represents a page.
    2. Each inner array inside the outer array has 6 strings, the first representing the topic for the slide and the follow 5 strings representing the facts in the trascript that pertain to the topic. The very first inner array represents the title page, so it only has one string for the title.
    3. The array should be the only thing you print in response to this prompt, in python form.
    4. Do not label any of the slides with any extra character outside the array. *Important*, only paste the array itself
    5. ALWAYS include a title for the title page FIRST, make it the title of the lecture or total theme of the presentation
    6. DO NOT include any image or spots to include images or new line "\n" characters
    7. NEVER use / or // characters
    8. Also keep the title in its own array layer inside the outer array, like it is a slide title but on its own
    9. Make sure you are generating long, detailed, and specific bullet points for a total of 20 slides. This is at the college level.
    9. Triple check that you are completely closing every array component that has square brackets. Make sure "[] was never closed" error never happens

""",
    'stream': True  # stream allows us to read the response incrementally
}

# === Call LLaMA and collect the LaTeX code ===
print("Generating Powerpoint blueprint w/ Llama 3.2")
response = requests.post(url, json=data, stream=True)  # POST request to the local model

if response.status_code == 200:
    pptx_code = ""  # buffer to hold the output
    for line in response.iter_lines():  # iterate over streamed lines
        if line:
            decoded_line = line.decode("utf-8")  # decode byte string
            result = json.loads(decoded_line)  # parse JSON from LLaMA response
            chunk = result.get("response", "")  # get response text
            print(chunk, end="", flush=True)  # live output to console
            pptx_code += chunk  # accumulate result
    print("\n\nDone generating pppt. Now compiling to Presentation")
else:
    print("Error:", response.status_code, response.text)
    pptx_code = None

# Convert the generated string output into a Python array (safe evaluation)
data_array = ast.literal_eval(pptx_code)

# Load the PowerPoint presentation using a template
template = "TemplatePro.pptx"
prs = Presentation(template)

# Access and modify the title slide
title_slide = prs.slides[0]
print(title_slide)
title_slide.shapes.title  # get title shape object
titleshape = title_slide.shapes.title
titleshape.text = data_array[0][0]  # set title to the first string in first sub-array

# Loop through the rest of the slides and populate them
count = 1
for item in data_array[1:]:
    if(count >20):  # safety guard to limit to 20 slides
        break
    textslide = prs.slides[count]  # get slide at index
    slidetitle = textslide.shapes.title  # get title placeholder
    slidetitle.text = item[0]  # set slide title
    bodyshape = textslide.shapes.placeholders[1]  # get body placeholder
    tf = bodyshape.text_frame  # access text frame for bullet points
    for bullet in item[1:]:  # loop through the 5 bullet points
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1  # indent level 1 for bullets
    count +=1

# Save the finished PowerPoint presentation
prs.save("test.pptx")

